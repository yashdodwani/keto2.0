from flask import Flask, request, jsonify
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
from threading import Thread
import pyttsx3
import html
from bs4 import BeautifulSoup
from flask_cors import CORS 
import re
import json
from langchain_community.llms import GPT4All  
from langchain_groq import ChatGroq
import os
from dotenv import load_dotenv
load_dotenv()
from pymongo import MongoClient
import PyPDF2
import google.generativeai as genai
import io
from google.generativeai import GenerativeModel
import jwt
from functools import wraps
from werkzeug.utils import secure_filename
import logging
from bson.objectid import ObjectId
from langchain_core.prompts import ChatPromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
#from langchain_community.indexes import VectorstoreIndexCreator  # Commented out due to ImportError
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.chains import RetrievalQA
from io import BytesIO
from PyPDF2 import PdfReader  
from langchain.schema import Document  
import chromadb
from sentence_transformers import SentenceTransformer
import google.generativeai as genai
#from langchain.document_loaders import PyPDFLoader
from pptx import Presentation
import redis
redis_client = redis.StrictRedis(host='localhost', port=6379, db=0, decode_responses=True)

app = Flask(__name__)
SECRET_KEY = os.getenv("SECRET_KEY", "changeme")
mongo_client = MongoClient("mongodb://localhost:27017/quicklearnai") 
db = mongo_client["quicklearnai"]
topics_collection = db["statistics"]

CORS(app, resources={
    r"/*": {
        "origins": ["http://localhost:5173", "http://localhost:3000", "http://localhost:3001"],
        "methods": ["GET", "POST", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization"],
        "supports_credentials": True
    }
})

formatter = TextFormatter()

google_api_key = os.getenv("GOOGLE_API_KEY")
if google_api_key:
    genai.configure(api_key=google_api_key)
gemini_model = genai.GenerativeModel('gemini-2.0-flash')  # Use the correct model name

groq_model = ChatGroq(
    model="llama-3.3-70b-specdec",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY")
)

def get_and_enhance_transcript(youtube_url, model_type='gemini'):
    try:
        video_id = youtube_url.split('v=')[-1]
        transcript = None
        language = None

        # Fetch transcript
        for lang in ['hi', 'en']:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=[lang])
                language = lang
                break
            except:
                continue

        if not transcript:
            return None, None

        formatted_transcript = "\n".join([entry['text'] for entry in transcript])

        # Enhanced transcript prompt
        prompt = f"""
        Act as a transcript cleaner. Generate a new transcript with the same context and content as the given transcript.
        If there’s a revision portion, differentiate it from the actual transcript.
        Output in sentences line by line. If the transcript lacks educational content, return 'Fake transcript'.
        Transcript: {formatted_transcript}
        """

        if model_type.lower() == 'chatgroq':
            response = groq_model.invoke(prompt)
            enhanced_transcript = response.content if hasattr(response, 'content') else str(response)
        else:  # Default to gemini
            google_api_key = os.getenv("GOOGLE_API_KEY")
            if google_api_key:
                genai.configure(api_key=google_api_key)
            gemini_model = genai.GenerativeModel('gemini-2.0-flash') 
            response = gemini_model.generate_content(prompt)
            enhanced_transcript = response.text if hasattr(response, 'text') else str(response)
        
        return enhanced_transcript, language
    except Exception as e:
        print(f"Error in get_and_enhance_transcript: {str(e)}")
        return None, None

def generate_summary_and_quiz(transcript, num_questions, language, difficulty, model_type='gemini'):
    try:
        if 'Fake transcript' in transcript:
            return {"summary": {}, "questions": {difficulty: []}}

        prompt = f"""
        Summarize the following transcript by identifying the key topics covered, and provide a detailed summary of each topic in 6-7 sentences.
        Each topic should be labeled clearly as "Topic X", where X is the topic name. Provide the full summary for each topic in English, even if the transcript is in a different language.
        Strictly ensure that possessives (e.g., John's book) and contractions (e.g., don't) use apostrophes (') instead of quotation marks (" or "  ").

        If the transcript contains 'Fake Transcript', do not generate any quiz or summary.

        After the summary, give the name of the topic on which the transcript was all about in a maximum of 2 to 3 words.
        After summarizing, create a quiz with {num_questions} multiple-choice questions in English, based on the transcript content.
        Only generate {difficulty} difficulty questions. Format the output in JSON format as follows, just give the JSON as output, nothing before it:

        {{
            "summary": {{
                "topic1": "value1",
                "topic2": "value2",
                "topic3": "value3"
            }},
            "questions": {{
                "{difficulty}": [
                    {{
                        "question": "What is the capital of France?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Paris"
                    }},
                    {{
                        "question": "What is the capital of Germany?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Berlin"
                    }}
                ]
            }}
        }}

        Transcript: {transcript}
        """

        if model_type.lower() == 'chatgroq':
            response = groq_model.invoke(prompt)
            response_content = response.content if hasattr(response, 'content') else str(response)
        else:  # Default to gemini
            response = gemini_model.generate_content(prompt)
            response_content = response.text if hasattr(response, 'text') else str(response)

        # Extract JSON from response
        json_match = re.search(r'\{.*\}', response_content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            try:
                return json.loads(json_str)
            except json.JSONDecodeError as e:
                print(f"JSONDecodeError: {e}, Raw response: {response_content}")
                return None
        else:
            print(f"No valid JSON found in response: {response_content}")
            return None
    except Exception as e:
        print(f"Error in generate_summary_and_quiz: {str(e)}")
        return None

@app.route('/quiz', methods=['POST', 'OPTIONS'])
def quiz():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    data = request.json
    youtube_link = data.get('link')
    num_questions = int(data.get('qno', 5))  # Default to 5 if not provided
    difficulty = data.get('difficulty', 'medium')  # Default to medium
    model_type = data.get('model', 'chatgroq')  # Default to gemini, can be 'chatgroq' or 'gemini'

    if not youtube_link:
        return jsonify({"error": "No YouTube URL provided"}), 400

    transcript, language = get_and_enhance_transcript(youtube_link, model_type)
    if not transcript:
        return jsonify({"error": "Failed to fetch transcript"}), 404

    summary_and_quiz = generate_summary_and_quiz(transcript, num_questions, language, difficulty, model_type)
    if summary_and_quiz:
        return jsonify(summary_and_quiz)
    else:
        return jsonify({"error": "Failed to generate quiz"}), 500

# recommendation
def validate_token_middleware():
    def middleware(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            auth_header = request.headers.get("Authorization")
            token = auth_header.split("Bearer ")[-1] if auth_header and "Bearer " in auth_header else None
            
            if not token:
                return jsonify({"message": "Unauthorized: No token provided"}), 401
            
            try:
                # Decoding the token using the correct jwt.decode()
                decoded = jwt.decode(token, SECRET_KEY, algorithms=["HS256"])
                request.user_id = decoded.get("id")
                request.user_role = decoded.get("role")  # Optional
                
                return func(*args, **kwargs)
            except jwt.ExpiredSignatureError:
                return jsonify({"message": "Unauthorized: Token has expired"}), 401
            except jwt.InvalidTokenError as e:
                print(f"Token decoding error: {e}")
                return jsonify({"message": "Unauthorized: Invalid token"}), 401
        
        return wrapper
    return middleware


from langchain.prompts import PromptTemplate
@app.route('/chat_trans', methods=['POST', 'OPTIONS'])
def chat_with_transcript():
    """Handle chat requests with YouTube transcript context"""
    if request.method == 'OPTIONS':
        return '', 204  # Handle CORS preflight request

    try:
        data = request.json
        if not data:
            return jsonify({'error': 'No JSON data provided'}), 400

        youtube_link = data.get('link')
        model_type = data.get('model', 'chatgroq')  # Default to chatgroq
        question = data.get('question')

        if not youtube_link:
            return jsonify({'error': 'Missing YouTube link'}), 400

        # Get and enhance transcript
        transcript, language = get_and_enhance_transcript(youtube_link, model_type)
        
        if "Error" in transcript:
            return jsonify({'error': transcript}), 400

        # If no question provided, just return the transcript
        if not question:
            return jsonify({
                'transcript': transcript,
                'language': language,
                'status': 'success'
            })

        # Process question with transcript context
        prompt_template = PromptTemplate(
            input_variables=["transcript", "question"],
            template="""Given the following YouTube video transcript:
            {transcript}
            
            Please answer this question based on the transcript content:
            {question}"""
        )

        formatted_prompt = prompt_template.format(
            transcript=transcript,
            question=question
        )

        # Get response from Groq
        response = groq_model.invoke(formatted_prompt)

        return jsonify({
            'answer': response.content,
            # 'transcript': transcript,
            # 'language': language,
            'status': 'success'
        })

    except Exception as e:
        return jsonify({'error': f'An error occurred: {str(e)}'}), 500

# Function to interact with LLaMA API
def llama_generate_recommendations(prompt):
    try:
        llm = ChatGroq(
            model="llama-3.3-70b-specdec",
            temperature=0,
            groq_api_key=os.getenv("GROQ_API_KEY")
        )
        
        response = llm.invoke(prompt)
        
        if hasattr(response, 'content'):
            return response.content
        else:
            return "Error: No content in response"
    except Exception as e:
        return f"Error connecting to Groq API: {e}"
 
 
 
import json

@app.route('/getonly', methods=['GET'])
@validate_token_middleware()
def get_recommendations():
    user_id = request.user_id  # Extract user ID from the token
    
    try:
        # Fetch user statistics from Redis
        statistics = redis_client.hget(f"student:{user_id}", "statistics")
        
        if not statistics:
            return jsonify({"message": "No statistics found for the provided user."}), 404
        
        # Convert JSON string to Python dictionary
        topics_data = json.loads(statistics)

        if not topics_data:
            return jsonify({"message": "No topics found for the provided user."}), 404

        # Extract only topic names
        topics_list = list(topics_data.keys())

        # Format recommendations prompt
        prompt = f"""
        Act as an intelligent recommendation generator. Based on the topics provided, generate a structured JSON response 
        with an overview, recommendations, and five YouTube video URLs for each topic. Ensure the output is in strict JSON 
        format without markdown or extra formatting. Use the following JSON structure:
        {{
            "topics": {{
                "<topic_name>": {{
                    "overview": "<brief overview>",
                    "recommendations": "<recommended steps to learn>",
                    "youtube_links": [
                        "<video_link_1>",
                        "<video_link_2>",
                        "<video_link_3>",
                        "<video_link_4>",
                        "<video_link_5>"
                    ]
                }}
            }}
        }}

        The topics are: {', '.join(topics_list)}
        """

        # Generate recommendations
        recommendations_raw = llama_generate_recommendations(prompt)

        # Ensure the response is valid JSON
        try:
            recommendations = json.loads(recommendations_raw)
        except json.JSONDecodeError:
            return jsonify({"message": "Failed to parse AI response as JSON", "raw_response": recommendations_raw}), 500

        return jsonify({
            "message": "Recommendations generated successfully",
            "recommendations": recommendations["topics"]  # Extract only relevant content
        }), 200

    except Exception as e:
        print("Error:", str(e))
        return jsonify({"message": f"An error occurred: {str(e)}"}), 500


import faiss 
from sentence_transformers import SentenceTransformer
from huggingface_hub import login
groq_api_key = os.getenv("GROQ_API_KEY")
groq_model_name = "llama3-8b-8192"
login(token=os.getenv("HUGGINGFACE_TOKEN")) 

groq_chat = ChatGroq(
    groq_api_key=groq_api_key,
    model_name=groq_model_name,
)


# Define the Groq system prompt
groq_sys_prompt = ChatPromptTemplate.from_template(
    "You are very smart at everything, you always give the best, the most accurate and most precise answers. "
    "Answer the following questions: {user_prompt}. Add more information as per your knowledge so that user can get proper knowledge, but make sure information is correct"
)
import threading
import time


embedding_model = SentenceTransformer('multi-qa-mpnet-base-cos-v1')  # Pre-trained model for embeddings
dimension = embedding_model.get_sentence_embedding_dimension()
faiss_index = faiss.IndexFlatL2(dimension) 
metadata_store = {}
pdf_storage = {}

def store_in_faiss(filename, text):
    chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
    embeddings = embedding_model.encode(chunks)
    faiss_index.add(embeddings)  
    metadata_store.update({i: filename for i in range(len(metadata_store), len(metadata_store) + len(chunks))})
genai.configure(api_key=os.getenv("GENAI_API_KEY"))
model = SentenceTransformer("all-MiniLM-L6-v2")
chroma_client = chromadb.PersistentClient(path="./chroma_db")
collection = chroma_client.get_or_create_collection(name="pdf_documents")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextToSpeechManager:
    def __init__(self):
        self.lock = threading.Lock()
    
    def speak(self, text):
        try:
            with self.lock:  # Ensure only one speech operation happens at a time
                engine = None
                try:
                    engine = pyttsx3.init()
                    engine.setProperty('rate', 150)
                    engine.setProperty('volume', 1.0)
                    engine.say(text)
                    engine.runAndWait() 
                    # print("spoke")
                    engine.startLoop(False)  # Start the event loop without blocking
                    engine.iterate()  # Process queued commands
                    engine.endLoop()  # End the event loop
                    logger.info("Speech completed successfully")
                finally:
                    if engine:
                        try:
                            engine.stop()
                        except:
                            pass
                        del engine
        except Exception as e:
            logger.error(f"Text-to-speech error: {str(e)}")
            
    def start_speaking(self, text):
        """Start a new thread for speaking"""
        thread = Thread(target=self.speak, args=(text,))
        thread.daemon = True  # Make thread daemon so it doesn't block program exit
        thread.start()
        return thread

# Create a global instance of the TTS manager
tts_manager = TextToSpeechManager()

def clean_response(text):
    """Clean and format the LLM response."""
    text = BeautifulSoup(text, "html.parser").get_text()
    text = html.unescape(text)
    text = re.sub(r'\n\s*\n', '\n\n', text)
    text = re.sub(r'[^\w\s.,!?-]', '', text)
    text = ' '.join(text.split())
    return text

def speak_text(text):
    """Convert text to speech using the TTS manager."""
    tts_manager.speak(text)

def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    return " ".join(page.extract_text() for page in reader.pages if page.extract_text())

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
    return " ".join(text)

@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400
    
    file_ext = os.path.splitext(file.filename)[-1].lower()
    file_path = os.path.join("./uploads", file.filename)
    os.makedirs("./uploads", exist_ok=True)
    file.save(file_path)
    
    try:
        if file_ext == ".pdf":
            content = extract_text_from_pdf(file_path)
        elif file_ext == ".pptx":
            content = extract_text_from_pptx(file_path)
        else:
            return jsonify({"error": "Unsupported file format. Only PDF and PPTX are allowed."}), 400
        
        existing_ids = collection.get()["ids"]
        if existing_ids:
            collection.delete(ids=existing_ids)
        embedding = model.encode(content).tolist()
        collection.add(documents=[content], embeddings=[embedding], ids=[file.filename])
        
        return jsonify({"message": "File uploaded and processed successfully."}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/test-audio", methods=["GET"])
def test_audio():
    try:
        test_text = "This is a test of the text to speech system"
        logger.info("Testing text-to-speech with test message")
        
        # Start speech in a new thread
        speech_thread = tts_manager.start_speaking(test_text)
        
        return jsonify({
            "message": "Audio test initiated",
            "test_text": test_text,
            "status": "Speech initiated"
        })
    except Exception as e:
        logger.error(f"Audio test failed: {str(e)}")
        return jsonify({
            "error": "Audio test failed",
            "details": str(e)
        }), 500

@app.route("/query", methods=["POST"])
def query_file():
    try:
        data = request.get_json()
        query = data.get("query", "")
        
        logger.info(f"Received query: {query}")
        
        query_embedding = model.encode(query).tolist()
        results = collection.query(query_embeddings=[query_embedding], n_results=3)
        retrieved_texts = "\n".join(results["documents"][0])
        
        prompt = f"""
        Based on the following context, please provide a clear and concise answer to the question.
        If the answer cannot be found in the context, please say so.
        
        Context: {retrieved_texts}
        
        Question: {query}
        """
        
        response = genai.GenerativeModel("gemini-1.5-flash").generate_content(prompt)
        cleaned_response = clean_response(response.text)
        
        # Add a small delay before starting new speech
        time.sleep(0.1)  # 100ms delay
        
        speech_thread = tts_manager.start_speaking(cleaned_response)
        
        return jsonify({
            "answer": cleaned_response,
            "voice_enabled": True,
            "status": "Speech initiated"
        })
        
    except Exception as e:
        error_message = f"Error processing query: {str(e)}"
        logger.error(error_message)
        return jsonify({
            "error": error_message,
            "answer": "I apologize, but I encountered an error while processing your query. Please try again.",
            "voice_enabled": False
        }), 500


# # Configure text-to-speech settings (optional)
# @app.route("/configure-voice", methods=["POST"])
# def configure_voice():
#     try:
#         data = request.get_json()
#         rate = data.get("rate", 110)  # Default speaking rate
#         volume = data.get("volume", 1.0)  # Default volume
#         voice_id = data.get("voice_id")  # Voice identifier
        
#         engine.setProperty('rate', rate)
#         engine.setProperty('volume', volume)
        
#         if voice_id:
#             voices = engine.getProperty('voices')
#             for voice in voices:
#                 if voice.id == voice_id:
#                     engine.setProperty('voice', voice.id)
#                     break
        
#         return jsonify({"message": "Voice settings updated successfully"}), 200
#     except Exception as e:
#         return jsonify({"error": f"Error configuring voice: {str(e)}"}), 500

# Configure text-to-speech settings (optional)
# @app.route("/configure-voice", methods=["POST"])
# def configure_voice():
#     try:
#         data = request.get_json()
#         rate = data.get("rate", 110)  # Default speaking rate
#         volume = data.get("volume", 1.0)  # Default volume
#         voice_id = data.get("voice_id")  # Voice identifier
        
#         engine.setProperty('rate', rate)
#         engine.setProperty('volume', volume)
        
#         if voice_id:
#             voices = engine.getProperty('voices')
#             for voice in voices:
#                 if voice.id == voice_id:
#                     engine.setProperty('voice', voice.id)
#                     break
        
#         return jsonify({"message": "Voice settings updated successfully"}), 200
#     except Exception as e:
#         return jsonify({"error": f"Error configuring voice: {str(e)}"}), 500
#         return jsonify({"error": f"Error configuring voice: {str(e)}"}), 500
# # MindMap

def fetch_youtube_transcript(video_url):
    try:
        video_id = video_url.split("v=")[-1]
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'hi'])
        return " ".join([entry["text"] for entry in transcript])  # Clean transcript
    except Exception as e:
        return {"error": f"Error fetching transcript: {str(e)}"}

def generate_mind_map(content):
    prompt = f"""
    Extract key concepts from the following text and structure them into a JSON-based mind map.
    Organize it into: "Topic" -> "Subtopics" -> "Details".

    Text: {content}

    Output **ONLY** valid JSON in this format (no extra text, no explanations):
    {{
        "topic": "Main Topic",
        "subtopics": [
            {{"name": "Subtopic 1", "details": ["Detail 1", "Detail 2"]}},
            {{"name": "Subtopic 2", "details": ["Detail 3", "Detail 4"]}}
        ]
    }}
    """

    llm = ChatGroq(
        model="llama-3.3-70b-specdec",
        temperature=0,
        groq_api_key=os.getenv("GROQ_API_KEY"),
    )      

    response = llm.invoke(prompt)

    # Ensure response is a string
    raw_json = response.content.strip() if hasattr(response, "content") else str(response)

    # Remove unwanted formatting (like triple backticks and newlines)
    cleaned_json_str = raw_json.replace("```json", "").replace("```", "").replace("\n", "").strip()

    # Convert to valid JSON
    try:
        return json.loads(cleaned_json_str)
    except json.JSONDecodeError:
        return {"error": f"Invalid JSON response: {cleaned_json_str}"}

@app.route("/generate_mind_map", methods=['GET'])
def generate_mind_map_endpoint():
    # print("✅ Endpoint called!")  # Debugging
    video_url = request.args.get('video_url')

    if not video_url:
        return jsonify({"error": "No video URL provided"}), 400

    transcript = fetch_youtube_transcript(video_url)
    if isinstance(transcript, dict) and "error" in transcript:
        return jsonify(transcript), 400

    mind_map = generate_mind_map(transcript)
   
    return jsonify(mind_map)

llm = ChatGroq(
    model="llama-3.3-70b-specdec",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY"),
)

def generate_quiz(topic: str, num_questions: int, difficulty: str):
    """Generate a quiz based on the given topic."""
    prompt = f"""
    Create a quiz on the topic: "{topic}". Generate {num_questions} multiple-choice questions.
    The questions should be of {difficulty} difficulty.
    Format the output strictly in JSON format as follows:
    
    {{
       
        "questions": {{
            "{difficulty}": [
                {{
                    "question": "What is ...?",
                    "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
                    "answer": "Option 1"
                }}
            ]
        }}
    }}
    """
    response = llm.invoke(prompt)
    return response.content if hasattr(response, 'content') else response.text

@app.route("/llm_quiz", methods=["POST"])
def quiz_endpoint():
    data = request.json
    topic = data.get("topic")
    num_questions = data.get("num_questions")
    difficulty = data.get("difficulty")
    
    if not topic:
        return jsonify({"error": "Topic is required"}), 400
    
    try:
        response_content = generate_quiz(topic, num_questions, difficulty)

        
        try:
            result = json.loads(response_content)
        except json.JSONDecodeError:
            json_start = response_content.find('{')
            json_end = response_content.rfind('}') + 1
            if json_start != -1 and json_end != -1:
                json_str = response_content[json_start:json_end]
                result = json.loads(json_str)
            else:
                return jsonify({"error": "Could not parse JSON from response"}), 500
        
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500




llm = ChatGroq(
    model="llama-3.3-70b-specdec",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY")
)

def generate_quiz(topic: str, num_questions: int, difficulty: str):
    """Generate a quiz based on the given topic."""
    prompt = f"""
    Create a quiz on the topic: "{topic}". Generate {num_questions} multiple-choice questions.
    The questions should be of {difficulty} difficulty.
    Format the output strictly in JSON format as follows:
    
    {{
       
        "questions": {{
            "{difficulty}": [
                {{
                    "question": "What is ...?",
                    "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
                    "answer": "Option 1"
                }}
            ]
        }}
    }}
    """
    response = llm.invoke(prompt)
    return response.content if hasattr(response, 'content') else response.text


@app.route('/', methods=['GET'])
def health():
    return jsonify({"status": "ok"}) 

if __name__ == '__main__':
    app.run(debug=True, port=5001)
    